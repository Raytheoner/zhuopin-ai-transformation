"""文档读取器 — docx/pdf → DocumentSections。

接口与未来平台 shared_tools/doc_parser.py 同构（rule-of-three触发时零改平移）。
"""
from __future__ import annotations

import re
from dataclasses import dataclass, field
from pathlib import Path

# D1-D8 段落标头的识别模式（宽松匹配多种格式）
_SECTION_PATTERNS = [
    # "D1:", "D1：", "D1 团队成员"
    re.compile(r"^[Dd][1-8]\s*[:：\s]", re.MULTILINE),
    # "第一步：", "第二步："
    re.compile(r"^第[一二三四五六七八]步\s*[:：]", re.MULTILINE),
    # "1.", "2." — only when at line start with space after dot
    re.compile(r"^\d\.\s+\S", re.MULTILINE),
    # "Step 1:", "Step1:"
    re.compile(r"^Step\s*\d\s*[:：]", re.MULTILINE | re.IGNORECASE),
]

_D_NUM_RE = re.compile(r"[Dd]([1-8])")


@dataclass
class DocumentSections:
    """解析后的文档结构。"""
    full_text: str
    # 按 D 编号索引的段落内容，key="D1"..."D8"
    sections: dict[str, str] = field(default_factory=dict)
    # 文档的前 500 字（标题/摘要区）
    header_text: str = ""
    source_path: str = ""


def read(path: Path | str) -> DocumentSections:
    """统一入口，按扩展名路由到 docx/pdf 读取器。"""
    p = Path(path)
    suffix = p.suffix.lower()
    if suffix in (".docx", ".doc"):
        return _read_docx(p)
    elif suffix == ".pdf":
        return _read_pdf(p)
    elif suffix in (".pptx", ".ppt"):
        return _read_pptx(p)
    else:
        raise ValueError(f"不支持的文件格式：{suffix}（支持 .docx/.pdf/.pptx）")


def _read_pptx(path: Path) -> DocumentSections:
    """读取 pptx：逐幻灯片提取形状文本 + 表格单元格，按阅读顺序拼接。

    8D 报告常为一张幻灯片一个 D 步骤（标题=D1/D2…），正文在文本框或表格里。
    """
    try:
        from pptx import Presentation  # type: ignore
    except ImportError as e:
        raise ImportError("需要 python-pptx：pip install python-pptx") from e

    prs = Presentation(str(path))
    parts: list[str] = []
    for slide in prs.slides:
        shapes = sorted(slide.shapes, key=lambda s: (s.top or 0, s.left or 0))
        for shape in shapes:
            parts.extend(_pptx_shape_text(shape))
    full_text = "\n".join(p for p in parts if p.strip())
    return _parse_sections(full_text, str(path))


def _pptx_shape_text(shape) -> list[str]:
    """从单个形状取文本：文本框逐段、表格逐行、组合形状递归。"""
    out: list[str] = []
    try:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = "".join(run.text for run in para.runs).strip()
                if t:
                    out.append(t)
        if shape.has_table:
            for row in shape.table.rows:
                cells = [c.text.strip() for c in row.cells]
                line = " | ".join(c for c in cells if c)
                if line:
                    out.append(line)
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            for sub in sorted(shape.shapes, key=lambda s: (s.top or 0, s.left or 0)):
                out.extend(_pptx_shape_text(sub))
    except Exception:
        pass  # 单个形状解析失败不影响整篇
    return out


def _read_docx(path: Path) -> DocumentSections:
    try:
        from docx import Document  # type: ignore
    except ImportError as e:
        raise ImportError("需要 python-docx：pip install python-docx") from e

    doc = Document(str(path))
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    full_text = "\n".join(paragraphs)
    return _parse_sections(full_text, str(path))


def _read_pdf(path: Path) -> DocumentSections:
    try:
        import pdfplumber  # type: ignore
    except ImportError as e:
        raise ImportError("需要 pdfplumber：pip install pdfplumber") from e

    lines = []
    with pdfplumber.open(str(path)) as pdf:
        for page in pdf.pages:
            text = page.extract_text() or ""
            lines.extend(text.splitlines())
    full_text = "\n".join(line.strip() for line in lines if line.strip())
    return _parse_sections(full_text, str(path))


# 8D 标准步骤标题 → D 编号（回退识别：报告标题只写中英文名、无「Dn.」前缀时）
_TITLE_TO_D = [
    (re.compile(r"^(?:D?1[.、\s]*)?团队(?:成员|组建)|^Team\s*member", re.I), "D1"),
    (re.compile(r"^(?:D?2[.、\s]*)?问题描述|^Problem\s*description", re.I), "D2"),
    (re.compile(r"^(?:D?3[.、\s]*)?(?:围堵|临时措施|遏制)|^Interim\s*Containment", re.I), "D3"),
    (re.compile(r"^(?:D?4[.、\s]*)?(?:根本原因|根因)|^Root\s*cause", re.I), "D4"),
    (re.compile(r"^(?:D?5[.、\s]*)?(?:纠正措施|永久措施|永久纠正)|^Permanent|^Corrective", re.I), "D5"),
    (re.compile(r"^(?:D?6[.、\s]*)?(?:效果验证|实施(?:及验证)?|验证措施)|^Implement|^Validat", re.I), "D6"),
    (re.compile(r"^(?:D?7[.、\s]*)?(?:预防措施|防止再发|举一反三)|^Prevent", re.I), "D7"),
    (re.compile(r"^(?:D?8[.、\s]*)?(?:结案|表彰|团队祝贺|恭喜)|^Closure|^Congratulat", re.I), "D8"),
]


def _title_section_key(line: str) -> str | None:
    """短标题行 → D 编号（仅对较短的标题行生效，避免正文中途误命中）。"""
    if len(line) > 40:
        return None
    for pat, key in _TITLE_TO_D:
        if pat.search(line):
            return key
    return None


def _parse_sections(full_text: str, source_path: str) -> DocumentSections:
    """从全文提取 D1–D8 段落。"""
    sections: dict[str, str] = {}
    header_text = full_text[:500]

    lines = full_text.splitlines()
    current_key: str | None = None
    current_lines: list[str] = []
    seen: set[str] = set()

    # 段落分隔符：冒号/空白/逗号/句点/顿号（pptx 8D 用「D1.问题描述」点号格式）
    sep = r"[:：\s，,.．、]"
    for line in lines:
        stripped = line.strip()
        d_match = _D_NUM_RE.match(stripped)
        # 宽松判定：行以 D\d 开头（含标点或空格紧跟）
        is_section_header = bool(
            d_match and re.match(rf"^[Dd][1-8]\s*{sep}", stripped)
        )
        key = f"D{d_match.group(1)}" if (is_section_header and d_match) else None
        rest_from = f"^[Dd][1-8]\\s*{sep}\\s*"
        # 回退：无「Dn.」前缀时，按 8D 标准步骤标题识别 D1/D2（部分报告标题只写中英文名）
        if key is None:
            tk = _title_section_key(stripped)
            if tk and tk not in seen:
                key = tk
                rest_from = None                     # 标题行整行即标头，正文从下一行起
        if key is not None:
            if current_key is not None:
                sections[current_key] = "\n".join(current_lines).strip()
            current_key = key
            seen.add(key)
            rest = re.sub(rest_from, "", stripped) if rest_from else ""
            current_lines = [rest] if rest else []
        else:
            if current_key is not None:
                current_lines.append(line)

    if current_key is not None:
        sections[current_key] = "\n".join(current_lines).strip()

    return DocumentSections(
        full_text=full_text,
        sections=sections,
        header_text=header_text,
        source_path=source_path,
    )
